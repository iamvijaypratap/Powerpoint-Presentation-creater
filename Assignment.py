from pptx import Presentation
from pptx.util import Inches
from wand.image import Image
import os
import pathlib
with Image(filename="nike_black.png") as img1:
    path = r'C:\Users\vijay pratap\Desktop\project'
    img1.resize(2500, 1000)
    img1.transparentize(0.0)
    # ---by using os module initiating a for loop to iterate images in directory
    for image in os.listdir(path):
        file_extension = pathlib.Path(image).suffix
        # --- if file extentsion is not .jpg for loop will skip these files and continue to another image files
        if file_extension != ".jpg":
            continue
        else:
            with Image(filename=image) as img2:
                img2.composite_channel("all_channels", img1, "dissolve", 0, 0)
                Name = image.split(".")
                img2.save(filename=f"Water_marked {Name[0]}.jpg")
#---declaring function which takes directory path as argument----
def create_ppt(directory):
    #---to remove old images which are not watermarked and there file extensions is not .jpg----
    files=[i for i in os.listdir(directory) if i.startswith("Water_") and i.endswith(".jpg")]
    X = Presentation()
    #----creating a dynamic layouts and slide of numbers of filtered images which is in directory----
    layouts=[]
    xslides=[]
    for j in range(len(files)):
        layouts.append("layouts"+str(j))
        xslides.append("slide"+str(j))
    #--- adding image files to slides---
    for x,y,file in zip(layouts,xslides,files):
        x = X.slide_layouts[1]
        y = X.slides.add_slide(x)

        y.shapes.title.text = "Creating a powerpoint using Python" 
        y.placeholders[1].text = "Created by Vijaypratap"
        top = Inches(3)
        left = Inches(1)
        height = Inches(4)
        pic = y.shapes.add_picture(file, left, top, height=height)
        X.save("presentation.pptx")
#--- path of directory---
directory=r'C:\Users\vijay pratap\Desktop\project'
create_ppt(directory)

